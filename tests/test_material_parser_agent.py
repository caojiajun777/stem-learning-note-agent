"""Tests for MaterialParserAgent integration with PDF/PPTX support.

Verifies that the agent correctly routes files by extension, writes
parsed/documents.json, and writes parsed/parse_warnings.md when warnings occur.
"""
from __future__ import annotations

import io
import json
from pathlib import Path

import pytest

from stem_learning_agent.core.config import RunConfig
from stem_learning_agent.core.workspace import CourseWorkspace
from stem_learning_agent.harness.orchestrator import Orchestrator


def _make_tiny_pptx_with_text(slide_texts: list[str]) -> bytes:
    """Generate a minimal PPTX with text on each slide."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for text in slide_texts:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        txBox.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_MINIMAL_2PAGE_PDF_BYTES = (
    b"%PDF-1.4\n"
    b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
    b"2 0 obj\n<< /Type /Pages /Kids [3 0 R 4 0 R] /Count 2 >>\nendobj\n"
    b"3 0 obj\n<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 5 0 R >> >> "
    b"/MediaBox [0 0 200 200] /Contents 6 0 R >>\nendobj\n"
    b"4 0 obj\n<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 5 0 R >> >> "
    b"/MediaBox [0 0 200 200] /Contents 7 0 R >>\nendobj\n"
    b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
    b"6 0 obj\n<< /Length 44 >>\nstream\nBT /F1 12 Tf 10 180 Td (Page 1) Tj ET\nendstream\nendobj\n"
    b"7 0 obj\n<< /Length 44 >>\nstream\nBT /F1 12 Tf 10 180 Td (Page 2) Tj ET\nendstream\nendobj\n"
    b"xref\n0 8\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n "
    b"\n0000000115 00000 n \n0000000262 00000 n \n0000000409 00000 n "
    b"\n0000000489 00000 n \n0000000582 00000 n \ntrailer\n<< /Size 8 /Root 1 0 R >>\n"
    b"startxref\n675\n%%EOF\n"
)


# ---------------------------------------------------------------------------
# 1. Agent routes files by extension and writes parsed/documents.json
# ---------------------------------------------------------------------------


def test_material_parser_agent_routes_by_extension(tmp_path: Path) -> None:
    """MaterialParserAgent must route .md, .pdf, .pptx correctly and write documents.json."""
    course = tmp_path / "course"
    raw = course / "raw"
    raw.mkdir(parents=True)

    # Create one of each supported type.
    (raw / "slides.md").write_text("# Slide 1\n\nContent.\n", encoding="utf-8")
    (raw / "textbook.pdf").write_bytes(_MINIMAL_2PAGE_PDF_BYTES)
    (raw / "examples.pptx").write_bytes(_make_tiny_pptx_with_text(["Example 1", "Example 2"]))

    cfg = RunConfig(course_path=course)
    orch = Orchestrator(cfg)
    orch.init()

    from stem_learning_agent.agents.material_parser_agent import MaterialParserAgent

    MaterialParserAgent().run(orch.ctx)

    # Check parsed/documents.json exists and has 3 documents.
    docs_path = CourseWorkspace(course).parsed_documents_path()
    assert docs_path.exists()
    raw_data = json.loads(docs_path.read_text(encoding="utf-8"))
    assert len(raw_data) == 3

    # Verify each document has chunks.
    md_doc = next((d for d in raw_data if d["material_id"] == "slides"), None)
    assert md_doc is not None
    assert len(md_doc["chunks"]) > 0

    pdf_doc = next((d for d in raw_data if d["material_id"] == "textbook"), None)
    assert pdf_doc is not None
    assert len(pdf_doc["chunks"]) == 2  # 2 pages

    pptx_doc = next((d for d in raw_data if d["material_id"] == "examples"), None)
    assert pptx_doc is not None
    assert len(pptx_doc["chunks"]) == 2  # 2 slides


# ---------------------------------------------------------------------------
# 2. Agent writes parsed/parse_warnings.md when warnings occur
# ---------------------------------------------------------------------------


def test_material_parser_agent_writes_warnings(tmp_path: Path) -> None:
    """When a parser emits warnings, they must land in parsed/parse_warnings.md."""
    course = tmp_path / "course"
    raw = course / "raw"
    raw.mkdir(parents=True)

    # Create a blank PDF (no text) to trigger a warning.
    try:
        from pypdf import PdfWriter
    except ImportError:
        pytest.skip("pypdf not installed")

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    (raw / "blank.pdf").write_bytes(buf.getvalue())

    cfg = RunConfig(course_path=course)
    orch = Orchestrator(cfg)
    orch.init()

    from stem_learning_agent.agents.material_parser_agent import MaterialParserAgent

    MaterialParserAgent().run(orch.ctx)

    warnings_path = CourseWorkspace(course).parse_warnings_path()
    assert warnings_path.exists()
    warnings_text = warnings_path.read_text(encoding="utf-8")
    assert "image-only" in warnings_text or "OCR" in warnings_text


# ---------------------------------------------------------------------------
# 3. Full pipeline still passes with sample course (Markdown-only baseline)
# ---------------------------------------------------------------------------


def test_full_pipeline_with_sample_course(sample_course_path: Path) -> None:
    """The existing sample course (Markdown-only) must still work end-to-end."""
    cfg = RunConfig(course_path=sample_course_path)
    orch = Orchestrator(cfg)
    orch.init()

    from stem_learning_agent.agents.curriculum_mapper_agent import CurriculumMapperAgent
    from stem_learning_agent.agents.material_parser_agent import MaterialParserAgent

    MaterialParserAgent().run(orch.ctx)
    CurriculumMapperAgent().run(orch.ctx)

    # Verify parsed/documents.json exists and has content.
    docs_path = CourseWorkspace(sample_course_path).parsed_documents_path()
    assert docs_path.exists()
    raw_data = json.loads(docs_path.read_text(encoding="utf-8"))
    assert len(raw_data) > 0
    # All chunks should have source_refs.
    for doc in raw_data:
        for chunk in doc["chunks"]:
            assert len(chunk["source_refs"]) > 0


# ---------------------------------------------------------------------------
# 4. Mixed course with .md + .pdf + .pptx
# ---------------------------------------------------------------------------


def test_mixed_course_md_pdf_pptx(tmp_path: Path) -> None:
    """A course with .md, .pdf, and .pptx files all parse correctly."""
    course = tmp_path / "mixed"
    raw = course / "raw"
    raw.mkdir(parents=True)

    (raw / "intro.md").write_text("# Introduction\n\nWelcome.\n", encoding="utf-8")
    (raw / "lecture.pdf").write_bytes(_MINIMAL_2PAGE_PDF_BYTES)
    (raw / "demo.pptx").write_bytes(_make_tiny_pptx_with_text(["Demo slide"]))

    cfg = RunConfig(course_path=course)
    orch = Orchestrator(cfg)
    orch.init()

    from stem_learning_agent.agents.material_parser_agent import MaterialParserAgent

    MaterialParserAgent().run(orch.ctx)

    docs_path = CourseWorkspace(course).parsed_documents_path()
    raw_data = json.loads(docs_path.read_text(encoding="utf-8"))
    assert len(raw_data) == 3

    # Verify each has chunks.
    for doc in raw_data:
        assert len(doc["chunks"]) > 0


if __name__ == "__main__":  # pragma: no cover
    pytest.main([__file__, "-v"])
