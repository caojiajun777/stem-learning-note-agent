"""parse_document: convert a raw material file into a ParsedDocument.

Supports:
- Markdown (.md, .markdown) — heading-aware chunking.
- Plain text (.txt) — single-chunk body.
- PDF (.pdf) — text-layer extraction, one chunk per page. No OCR.
- PPTX (.pptx) — text extraction from slide shapes, one chunk per slide. No OCR.

Unsupported extensions produce a warning and an empty ParsedDocument so the
pipeline can continue without crashing.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Optional

from ..core.errors import ToolError
from ..core.logging import get_logger
from ..core.schemas import ParsedChunk, ParsedDocument, SourceRef
from ..harness.tool_base import Tool, ToolResult

log = get_logger(__name__)

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
_BULLET_RE = re.compile(r"^[-*]\s+(.+)$")


# ---------------------------------------------------------------------------
# Markdown parser (unchanged)
# ---------------------------------------------------------------------------


def _split_markdown(text: str, material_id: str) -> list[ParsedChunk]:
    chunks: list[ParsedChunk] = []
    current_heading: str | None = None
    buffer: list[str] = []
    buffer_start = 1

    def flush(end_line: int) -> None:
        if not buffer:
            return
        body = "\n".join(buffer).strip()
        if not body:
            return
        cid = f"{material_id}-c{len(chunks):03d}"
        chunks.append(
            ParsedChunk(
                id=cid,
                material_id=material_id,
                text=body,
                heading=current_heading,
                chunk_type="body",
                source_refs=[
                    SourceRef(
                        material_id=material_id,
                        line_start=buffer_start,
                        line_end=end_line,
                    )
                ],
                confidence=0.9,
            )
        )

    lines = text.splitlines()
    for i, line in enumerate(lines, start=1):
        m = _HEADING_RE.match(line)
        if m:
            flush(i - 1)
            current_heading = m.group(2).strip()
            buffer = []
            buffer_start = i + 1
            cid = f"{material_id}-h{len(chunks):03d}"
            chunks.append(
                ParsedChunk(
                    id=cid,
                    material_id=material_id,
                    text=current_heading,
                    heading=current_heading,
                    chunk_type="title",
                    source_refs=[
                        SourceRef(
                            material_id=material_id,
                            line_start=i,
                            line_end=i,
                        )
                    ],
                    confidence=0.95,
                )
            )
            continue
        if not buffer:
            buffer_start = i
        buffer.append(line)
    flush(len(lines))
    return chunks


# ---------------------------------------------------------------------------
# PDF parser (text-layer only, no OCR)
# ---------------------------------------------------------------------------


def _parse_pdf(
    path: Path, material_id: str
) -> tuple[list[ParsedChunk], str, list[str]]:
    """Extract text from a PDF, one chunk per page.

    Returns (chunks, extracted_text, warnings).
    If the PDF has no text layer or extraction fails, returns empty chunks
    and a clear warning. Does NOT crash the pipeline.
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        return (
            [],
            "",
            [
                f"parse_document: pypdf not installed; cannot parse '{path.name}'. "
                "Install via: pip install pypdf"
            ],
        )

    chunks: list[ParsedChunk] = []
    warnings: list[str] = []
    all_text_lines: list[str] = []

    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001 — pypdf can raise many exception types
        log.warning("parse_document: PDF open failed for %s: %s", path.name, exc)
        return (
            [],
            "",
            [
                f"parse_document: failed to open PDF '{path.name}': {type(exc).__name__}. "
                "File may be encrypted, corrupted, or unsupported."
            ],
        )

    if not reader.pages:
        return (
            [],
            "",
            [f"parse_document: PDF '{path.name}' has zero pages."],
        )

    for page_num, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            log.warning(
                "parse_document: text extraction failed for page %d of %s: %s",
                page_num,
                path.name,
                exc,
            )
            text = ""
            warnings.append(
                f"parse_document: page {page_num} of '{path.name}' text extraction failed: {type(exc).__name__}"
            )

        text = text.strip()
        all_text_lines.append(text)

        if not text:
            # Empty page — not an error, but note it.
            warnings.append(
                f"parse_document: page {page_num} of '{path.name}' has no extractable text "
                "(may be image-only or blank)."
            )
            continue

        cid = f"{material_id}-p{page_num:03d}"
        chunks.append(
            ParsedChunk(
                id=cid,
                material_id=material_id,
                text=text,
                chunk_type="body",
                source_refs=[
                    SourceRef(
                        material_id=material_id,
                        page=page_num,
                    )
                ],
                confidence=0.85,  # text-layer extraction is reliable but not perfect
            )
        )

    extracted_text = "\n\n".join(all_text_lines)
    if not chunks:
        warnings.append(
            f"parse_document: PDF '{path.name}' produced zero text chunks. "
            "The file may be image-only or scanned without OCR. "
            "OCR is not implemented in this MVP."
        )
    return chunks, extracted_text, warnings


# ---------------------------------------------------------------------------
# PPTX parser (text from slide shapes, no OCR)
# ---------------------------------------------------------------------------


def _parse_pptx(
    path: Path, material_id: str
) -> tuple[list[ParsedChunk], str, list[str]]:
    """Extract text from a PPTX, one chunk per slide.

    Returns (chunks, extracted_text, warnings).
    If a slide has no text, a warning is emitted but the pipeline continues.
    Does NOT crash on missing python-pptx or unreadable files.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return (
            [],
            "",
            [
                f"parse_document: python-pptx not installed; cannot parse '{path.name}'. "
                "Install via: pip install python-pptx"
            ],
        )

    chunks: list[ParsedChunk] = []
    warnings: list[str] = []
    all_text_lines: list[str] = []

    try:
        prs = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        log.warning("parse_document: PPTX open failed for %s: %s", path.name, exc)
        return (
            [],
            "",
            [
                f"parse_document: failed to open PPTX '{path.name}': {type(exc).__name__}. "
                "File may be corrupted or unsupported."
            ],
        )

    if not prs.slides:
        return (
            [],
            "",
            [f"parse_document: PPTX '{path.name}' has zero slides."],
        )

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        try:
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                t = shape.text.strip()
                if t:
                    text_parts.append(t)
        except Exception as exc:  # noqa: BLE001
            log.warning(
                "parse_document: shape iteration failed for slide %d of %s: %s",
                slide_num,
                path.name,
                exc,
            )
            warnings.append(
                f"parse_document: slide {slide_num} of '{path.name}' shape iteration failed: {type(exc).__name__}"
            )

        text = "\n".join(text_parts).strip()
        all_text_lines.append(text)

        if not text:
            warnings.append(
                f"parse_document: slide {slide_num} of '{path.name}' has no extractable text "
                "(may be image-only or blank)."
            )
            continue

        cid = f"{material_id}-s{slide_num:03d}"
        chunks.append(
            ParsedChunk(
                id=cid,
                material_id=material_id,
                text=text,
                chunk_type="body",
                source_refs=[
                    SourceRef(
                        material_id=material_id,
                        page=slide_num,  # slide number maps to page
                    )
                ],
                confidence=0.85,
            )
        )

    extracted_text = "\n\n".join(all_text_lines)
    if not chunks:
        warnings.append(
            f"parse_document: PPTX '{path.name}' produced zero text chunks. "
            "All slides may be image-only or blank. "
            "OCR is not implemented in this MVP."
        )
    return chunks, extracted_text, warnings


# ---------------------------------------------------------------------------
# Tool
# ---------------------------------------------------------------------------


class ParseDocumentTool(Tool):
    name = "parse_document"
    description = "Parse a raw course-material file into a ParsedDocument."

    def run(
        self, *, material_id: str, path: Path, material_type: str = "other"
    ) -> ToolResult:  # type: ignore[override]
        p = Path(path)
        warnings: list[str] = []
        if not p.exists():
            raise ToolError(f"parse_document: not found: {p}")

        suffix = p.suffix.lower()

        # Markdown / plain text (unchanged)
        if suffix in (".md", ".markdown", ".txt"):
            text = p.read_text(encoding="utf-8")
            if suffix == ".txt":
                chunks = [
                    ParsedChunk(
                        id=f"{material_id}-c000",
                        material_id=material_id,
                        text=text.strip(),
                        chunk_type="body",
                        source_refs=[
                            SourceRef(
                                material_id=material_id,
                                line_start=1,
                                line_end=text.count("\n") + 1,
                            )
                        ],
                        confidence=0.85,
                    )
                ]
            else:
                chunks = _split_markdown(text, material_id)
            doc = ParsedDocument(
                material_id=material_id,
                chunks=chunks,
                extracted_text=text,
                warnings=warnings,
            )
            return ToolResult(ok=True, data=doc, warnings=warnings)

        # PDF (text-layer only)
        if suffix == ".pdf":
            chunks, extracted_text, pdf_warnings = _parse_pdf(p, material_id)
            warnings.extend(pdf_warnings)
            doc = ParsedDocument(
                material_id=material_id,
                chunks=chunks,
                extracted_text=extracted_text,
                warnings=warnings,
            )
            return ToolResult(ok=True, data=doc, warnings=warnings)

        # PPTX (text from shapes only)
        if suffix == ".pptx":
            chunks, extracted_text, pptx_warnings = _parse_pptx(p, material_id)
            warnings.extend(pptx_warnings)
            doc = ParsedDocument(
                material_id=material_id,
                chunks=chunks,
                extracted_text=extracted_text,
                warnings=warnings,
            )
            return ToolResult(ok=True, data=doc, warnings=warnings)

        # Unsupported
        if suffix == ".docx":
            warnings.append(
                f"parse_document: '.docx' parsing is not implemented in MVP. "
                f"File '{p.name}' produced no chunks. See docs/tasks/02_document_parser.md."
            )
        else:
            warnings.append(
                f"parse_document: unsupported extension '{suffix}'. File '{p.name}' ignored."
            )
        return ToolResult(
            ok=True,
            data=ParsedDocument(
                material_id=material_id, extracted_text="", warnings=warnings
            ),
            warnings=warnings,
        )
